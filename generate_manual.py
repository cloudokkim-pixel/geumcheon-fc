"""
G.C.S.C 홈페이지 사용자 메뉴얼 PPT 자동 생성 스크립트

사전 설치:
    pip install playwright python-pptx pillow
    playwright install chromium

실행:
    python generate_manual.py                    # 로컬 개발 서버 (localhost:3000)
    python generate_manual.py --url https://geumcheon-fc.vercel.app   # 운영 서버
    python generate_manual.py --url http://localhost:3000 --out 메뉴얼.pptx
"""

import argparse
import io
import os
import sys
import time
from pathlib import Path

# ──────────────────────────────────────────────────────────────
# 페이지 목록 정의
# ──────────────────────────────────────────────────────────────
PAGES = [
    # (경로, 슬라이드 제목, 설명, 섹션)
    ("/",                  "메인 홈",        "클럽 소개·통계·파이프라인·갤러리 미리보기 등 사이트 전체 요약 페이지입니다.",                     "일반"),
    ("/club-introduction", "클럽 소개",      "서울금천축구클럽의 역사·철학·코칭스태프 등 클럽 정체성을 안내합니다.",                             "일반"),
    ("/program",           "육성 시스템",    "초등~성인 단계별 훈련 프로그램과 커리큘럼을 확인할 수 있습니다.",                                   "일반"),
    ("/application",       "선수 모집",      "전 연령 선수 모집 안내 및 지원 절차를 확인합니다.",                                                   "일반"),
    ("/gallery",           "갤러리",         "훈련·경기·행사 사진을 카테고리별로 탐색합니다.",                                                        "일반"),
    ("/support",           "후원 안내",      "클럽 후원 방법과 혜택을 안내합니다.",                                                                   "일반"),
    ("/support/apply",     "후원 신청",      "후원 신청 양식 페이지입니다.",                                                                          "일반"),
    ("/contact",           "문의",           "전화·이메일·SNS 등 클럽 연락처 및 위치 정보를 제공합니다.",                                             "일반"),
    # 관리자
    ("/admin/login",       "관리자 로그인",  "관리자 계정으로 로그인하는 페이지입니다. (일반 사용자 접근 불가)",                                       "관리자"),
    ("/admin",             "관리자 대시보드","로그인 후 접근하는 관리자 홈 화면입니다.",                                                               "관리자"),
    ("/admin/gallery",     "갤러리 관리",    "등록된 갤러리 목록을 확인하고 삭제할 수 있습니다.",                                                      "관리자"),
    ("/admin/gallery/new", "갤러리 등록",    "새 갤러리 항목(사진·제목·카테고리)을 등록하는 폼 페이지입니다.",                                         "관리자"),
]

# ──────────────────────────────────────────────────────────────
# 색상 / 디자인 상수
# ──────────────────────────────────────────────────────────────
C_BG        = "0D0D0D"   # 슬라이드 배경
C_ACCENT    = "CC2222"   # 강조 빨강
C_WHITE     = "F5F5F7"   # 밝은 흰색
C_GRAY      = "888890"   # 회색 텍스트
C_SECTION   = "1C1C1E"   # 섹션 구분 배경
C_ADMIN_ACC = "2255CC"   # 관리자 강조 파랑


def hex_rgb(h: str):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def capture_screenshots(base_url: str, pages: list, out_dir: Path) -> dict[str, bytes]:
    from playwright.sync_api import sync_playwright

    out_dir.mkdir(exist_ok=True)
    shots: dict[str, bytes] = {}

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        ctx = browser.new_context(
            viewport={"width": 1440, "height": 900},
            device_scale_factor=1,
        )
        page = ctx.new_page()

        # 관리자 페이지 미리 로그인 시도 (실패해도 계속 진행)
        try:
            page.goto(f"{base_url}/admin/login", timeout=15000, wait_until="networkidle")
            time.sleep(1)
        except Exception:
            pass

        for path, title, *_ in pages:
            url = f"{base_url}{path}"
            print(f"  캡처 중: {url}")
            try:
                page.goto(url, timeout=20000, wait_until="networkidle")
                page.wait_for_timeout(1200)
                # 고정 헤더 있으면 스크롤 약간 내렸다 올려 lazy 이미지 로드
                page.evaluate("window.scrollTo(0, 400)")
                page.wait_for_timeout(600)
                page.evaluate("window.scrollTo(0, 0)")
                page.wait_for_timeout(400)
                shot = page.screenshot(full_page=True)
                shots[path] = shot
                # 파일로도 저장
                (out_dir / f"{path.strip('/').replace('/', '_') or 'home'}.png").write_bytes(shot)
            except Exception as e:
                print(f"    ⚠ 실패 ({e}) — 빈 슬라이드로 대체")
                shots[path] = b""

        browser.close()

    return shots


def build_pptx(shots: dict[str, bytes], pages: list, out_path: Path):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree
    from PIL import Image

    SLIDE_W = Inches(13.33)
    SLIDE_H = Inches(7.5)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # completely blank

    def set_bg(slide, hex_color: str):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = hex_rgb(hex_color)

    def add_text(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
        from pptx.util import Pt, Inches
        txb = slide.shapes.add_textbox(left, top, width, height)
        txb.word_wrap = wrap
        tf = txb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = hex_rgb(color)
        run.font.name = "Malgun Gothic"
        return txb

    def add_rect(slide, left, top, width, height, hex_color: str, alpha=None):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_rgb(hex_color)
        shape.line.fill.background()
        return shape

    # ── 표지 슬라이드 ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, C_BG)

    # 빨간 왼쪽 바
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, C_ACCENT)

    add_text(slide, "G.C.S.C",
             Inches(0.5), Inches(1.8), Inches(12), Inches(0.8),
             font_size=14, color=C_GRAY)
    add_text(slide, "서울금천축구클럽",
             Inches(0.5), Inches(2.4), Inches(12), Inches(1.2),
             font_size=48, bold=True, color=C_WHITE)
    add_text(slide, "홈페이지 사용자 메뉴얼",
             Inches(0.5), Inches(3.4), Inches(12), Inches(0.8),
             font_size=28, bold=False, color=C_ACCENT)
    add_text(slide, "본 문서는 G.C.S.C 홈페이지의 주요 메뉴와 관리자 기능을 안내합니다.",
             Inches(0.5), Inches(4.4), Inches(10), Inches(0.6),
             font_size=14, color=C_GRAY)
    add_text(slide, "Since 2015  ·  geumcheon-fc.vercel.app",
             Inches(0.5), Inches(6.6), Inches(10), Inches(0.4),
             font_size=11, color=C_GRAY)

    # ── 목차 슬라이드 ─────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, C_BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, C_ACCENT)
    add_text(slide, "목  차", Inches(0.5), Inches(0.4), Inches(12), Inches(0.6),
             font_size=28, bold=True, color=C_WHITE)

    user_pages  = [(i + 1, p) for i, p in enumerate(pages) if p[3] == "일반"]
    admin_pages = [(i + 1, p) for i, p in enumerate(pages) if p[3] == "관리자"]
    base_no = 3  # 표지(1) + 목차(1) + 일반섹션(1) = 슬라이드 3부터 시작

    col_x = [Inches(0.6), Inches(6.8)]
    row_y_start = Inches(1.3)
    row_h = Inches(0.38)

    add_text(slide, "[ 일반 메뉴 ]", col_x[0], row_y_start - Inches(0.1), Inches(5.5), Inches(0.35),
             font_size=12, bold=True, color=C_ACCENT)
    for idx, (slide_no, (path, title, desc, _)) in enumerate(user_pages):
        y = row_y_start + Inches(0.3) + row_h * idx
        add_text(slide, f"{idx + 1:02d}.  {title}",
                 col_x[0], y, Inches(5.5), Inches(0.35),
                 font_size=13, color=C_WHITE)

    add_text(slide, "[ 관리자 메뉴 ]", col_x[1], row_y_start - Inches(0.1), Inches(5.5), Inches(0.35),
             font_size=12, bold=True, color=C_ADMIN_ACC)
    for idx, (slide_no, (path, title, desc, _)) in enumerate(admin_pages):
        y = row_y_start + Inches(0.3) + row_h * idx
        add_text(slide, f"{idx + 1:02d}.  {title}",
                 col_x[1], y, Inches(5.5), Inches(0.35),
                 font_size=13, color=C_WHITE)

    # ── 섹션 구분 슬라이드 생성 함수 ──────────────────────────
    def add_section_slide(label: str, sub: str, accent: str):
        sl = prs.slides.add_slide(blank_layout)
        set_bg(sl, C_SECTION)
        add_rect(sl, Inches(0), Inches(0), Inches(0.18), SLIDE_H, accent)
        add_text(sl, label, Inches(0.5), Inches(2.8), Inches(12), Inches(1.0),
                 font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, sub, Inches(0.5), Inches(3.9), Inches(12), Inches(0.6),
                 font_size=16, color=C_GRAY, align=PP_ALIGN.CENTER)

    # ── 일반 페이지 슬라이드 ─────────────────────────────────
    add_section_slide("일반 메뉴", "홈페이지 주요 페이지 안내", C_ACCENT)

    for path, title, desc, section in pages:
        if section != "일반":
            continue

        sl = prs.slides.add_slide(blank_layout)
        set_bg(sl, C_BG)

        # 상단 헤더 바
        add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.7), "1C1C1E")
        add_rect(sl, Inches(0), Inches(0), Inches(0.18), Inches(0.7), C_ACCENT)

        add_text(sl, title, Inches(0.35), Inches(0.1), Inches(8), Inches(0.5),
                 font_size=18, bold=True, color=C_WHITE)
        add_text(sl, path,  Inches(8.5), Inches(0.15), Inches(4.5), Inches(0.4),
                 font_size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)

        # 스크린샷 영역 (좌측 65%)
        shot_bytes = shots.get(path, b"")
        if shot_bytes:
            img_stream = io.BytesIO(shot_bytes)
            img = Image.open(img_stream)
            iw, ih = img.size
            # 비율 유지하며 최대 높이 6.3인치
            max_h = Inches(6.3)
            max_w = Inches(8.4)
            ratio = min(max_w / iw, max_h / ih)
            disp_w = int(iw * ratio)
            disp_h = int(ih * ratio)
            img_stream.seek(0)
            sl.shapes.add_picture(img_stream, Inches(0.25), Inches(0.85), disp_w, disp_h)

        # 우측 설명 패널
        panel_x = Inches(8.85)
        add_rect(sl, panel_x - Inches(0.1), Inches(0.75), Inches(4.35), Inches(6.5), "1C1C1E")

        add_text(sl, "페이지 안내", panel_x, Inches(0.9), Inches(4.0), Inches(0.4),
                 font_size=11, bold=True, color=C_ACCENT)
        add_text(sl, title, panel_x, Inches(1.25), Inches(4.0), Inches(0.7),
                 font_size=20, bold=True, color=C_WHITE)
        add_text(sl, desc, panel_x, Inches(2.0), Inches(4.0), Inches(2.0),
                 font_size=13, color=C_GRAY, wrap=True)
        add_text(sl, f"URL: {path}", panel_x, Inches(6.6), Inches(4.0), Inches(0.4),
                 font_size=10, color=C_GRAY)

    # ── 관리자 페이지 슬라이드 ───────────────────────────────
    add_section_slide("관리자 메뉴", "관리자 전용 기능 안내 (로그인 필요)", C_ADMIN_ACC)

    for path, title, desc, section in pages:
        if section != "관리자":
            continue

        sl = prs.slides.add_slide(blank_layout)
        set_bg(sl, C_BG)

        add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.7), "1C1C1E")
        add_rect(sl, Inches(0), Inches(0), Inches(0.18), Inches(0.7), C_ADMIN_ACC)

        add_text(sl, f"[관리자]  {title}", Inches(0.35), Inches(0.1), Inches(8), Inches(0.5),
                 font_size=18, bold=True, color=C_WHITE)
        add_text(sl, path, Inches(8.5), Inches(0.15), Inches(4.5), Inches(0.4),
                 font_size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)

        shot_bytes = shots.get(path, b"")
        if shot_bytes:
            img_stream = io.BytesIO(shot_bytes)
            img = Image.open(img_stream)
            iw, ih = img.size
            max_h = Inches(6.3)
            max_w = Inches(8.4)
            ratio = min(max_w / iw, max_h / ih)
            disp_w = int(iw * ratio)
            disp_h = int(ih * ratio)
            img_stream.seek(0)
            sl.shapes.add_picture(img_stream, Inches(0.25), Inches(0.85), disp_w, disp_h)

        panel_x = Inches(8.85)
        add_rect(sl, panel_x - Inches(0.1), Inches(0.75), Inches(4.35), Inches(6.5), "1C1C1E")
        add_text(sl, "관리자 기능", panel_x, Inches(0.9), Inches(4.0), Inches(0.4),
                 font_size=11, bold=True, color=C_ADMIN_ACC)
        add_text(sl, title, panel_x, Inches(1.25), Inches(4.0), Inches(0.7),
                 font_size=20, bold=True, color=C_WHITE)
        add_text(sl, desc, panel_x, Inches(2.0), Inches(4.0), Inches(2.0),
                 font_size=13, color=C_GRAY, wrap=True)
        add_text(sl, f"URL: {path}", panel_x, Inches(6.6), Inches(4.0), Inches(0.4),
                 font_size=10, color=C_GRAY)

    prs.save(str(out_path))
    print(f"\n[완료] 저장: {out_path}")


# ──────────────────────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(description="G.C.S.C 사이트 메뉴얼 PPT 생성")
    parser.add_argument("--url", default="http://localhost:3000",
                        help="사이트 기본 URL (기본값: http://localhost:3000)")
    parser.add_argument("--out", default="GCSC_사용자메뉴얼.pptx",
                        help="출력 파일명 (기본값: GCSC_사용자메뉴얼.pptx)")
    parser.add_argument("--shots-dir", default="manual_screenshots",
                        help="스크린샷 저장 폴더 (기본값: manual_screenshots)")
    parser.add_argument("--skip-capture", action="store_true",
                        help="이미 캡처된 스크린샷을 재사용하고 캡처를 건너뜀")
    args = parser.parse_args()

    base_url = args.url.rstrip("/")
    out_path = Path(args.out)
    shots_dir = Path(args.shots_dir)

    print(f"대상 URL : {base_url}")
    print(f"출력 파일: {out_path}")
    print(f"총 페이지: {len(PAGES)}개\n")

    if args.skip_capture and shots_dir.exists():
        print("── 저장된 스크린샷 불러오는 중 ──")
        shots: dict[str, bytes] = {}
        for path, *_ in PAGES:
            fname = f"{path.strip('/').replace('/', '_') or 'home'}.png"
            fpath = shots_dir / fname
            if fpath.exists():
                shots[path] = fpath.read_bytes()
                print(f"  로드: {fname}")
            else:
                print(f"  없음: {fname} (빈 슬라이드)")
                shots[path] = b""
    else:
        print("── 스크린샷 캡처 중 ──")
        shots = capture_screenshots(base_url, PAGES, shots_dir)

    print("\n── PPT 생성 중 ──")
    build_pptx(shots, PAGES, out_path)


if __name__ == "__main__":
    main()
